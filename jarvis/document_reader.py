"""jarvis.document_reader

Document reading and processing system for JARVIS.

Supports:
- PDF files (.pdf)
- Word documents (.docx, .doc)
- Text files (.txt, .md, .csv, .json)
- Excel files (.xlsx, .xls)
- PowerPoint (.pptx)
- Images with OCR (.png, .jpg, .jpeg)

Extracted content is saved to memory for future reference.
"""

from __future__ import annotations

import logging
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger(__name__)


class DocumentReader:
    """Read and extract text from various document formats."""
    
    SUPPORTED_EXTENSIONS = {
        '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm',
        '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt',
        '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif'
    }
    
    def __init__(self):
        self.extracted_cache: Dict[str, str] = {}
        
    def read_document(self, file_path: str, save_to_memory: bool = True) -> Dict[str, Any]:
        """Read any supported document and extract text.
        
        Args:
            file_path: Path to the document file
            save_to_memory: Whether to save extracted content to JARVIS memory
            
        Returns:
            Dict with extracted content, metadata, and status
        """
        try:
            path = Path(file_path)
            if not path.exists():
                return {
                    "success": False,
                    "error": f"File not found: {file_path}",
                    "content": ""
                }
            
            ext = path.suffix.lower()
            if ext not in self.SUPPORTED_EXTENSIONS:
                return {
                    "success": False,
                    "error": f"Unsupported file format: {ext}. Supported: {', '.join(self.SUPPORTED_EXTENSIONS)}",
                    "content": ""
                }
            
            # Extract based on file type
            if ext == '.pdf':
                content = self._read_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                content = self._read_word(file_path)
            elif ext in ['.xlsx', '.xls']:
                content = self._read_excel(file_path)
            elif ext in ['.pptx', '.ppt']:
                content = self._read_powerpoint(file_path)
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.gif']:
                content = self._read_image_ocr(file_path)
            else:
                # Text files
                content = self._read_text(file_path)
            
            if not content or content.strip() == "":
                logger.warning(f"Empty content extracted from {file_path}")
                return {
                    "success": False,
                    "error": f"Could not extract content from file (empty or corrupted). File type: {ext}, path: {file_path}",
                    "content": ""
                }
            
            # Truncate very long content for summary
            summary = content[:500] + "..." if len(content) > 500 else content
            
            result = {
                "success": True,
                "file_path": str(path.absolute()),
                "file_name": path.name,
                "file_type": ext,
                "file_size": path.stat().st_size,
                "content": content,
                "summary": summary,
                "content_length": len(content),
                "word_count": len(content.split()),
                "lines": content.count('\n') + 1
            }
            
            # Save to memory if requested (only if explicitly asked to learn/save)
            if save_to_memory:
                from jarvis.intent_classifier import should_save
                # Only save if user explicitly wants to save
                # The calling code should determine this based on intent
                pass  # Saving is now handled by intent classifier in main.py
            
            return result
            
        except Exception as e:
            logger.exception(f"Error reading document: {file_path}")
            return {
                "success": False,
                "error": f"Error reading document: {str(e)}",
                "content": ""
            }
    
    def read_document_clean(self, file_path: str, user_request: str = "") -> str:
        """Read document and return ONLY content, no metadata.
        
        Args:
            file_path: Path to document
            user_request: What user asked for (to extract specific parts)
            
        Returns:
            Clean content string only
        """
        result = self.read_document(file_path, save_to_memory=False)
        
        if not result.get("success"):
            return f"Error: {result.get('error', 'Could not read file')}"
        
        content = result.get("content", "")
        
        # If user asked for specific things, try to extract
        user_lower = user_request.lower()
        
        # Extract only code if requested
        if any(word in user_lower for word in ["code", "program", "script", "function"]):
            return self._extract_code_sections(content)
        
        # Extract only summary if requested
        if any(word in user_lower for word in ["summary", "summarize", "overview", "brief"]):
            return self._summarize_content(content)
        
        # Extract specific sections
        if any(word in user_lower for word in ["heading", "title", "section"]):
            return self._extract_headings(content)
        
        # Return full content without metadata prefix/suffix
        return content
    
    def _extract_code_sections(self, content: str) -> str:
        """Extract code blocks from content."""
        import re
        
        # Match code blocks
        code_patterns = [
            r'```[\w]*\n(.*?)```',
            r'`([^`]+)`',
        ]
        
        all_code = []
        for pattern in code_patterns:
            matches = re.findall(pattern, content, re.DOTALL)
            all_code.extend(matches)
        
        if all_code:
            return "\n\n--- Code Section ---\n\n".join(all_code)
        
        # If no code blocks found, return lines that look like code
        lines = content.split('\n')
        code_lines = []
        for line in lines:
            # Heuristic: lines with common code patterns
            if any(pattern in line for pattern in ['def ', 'class ', 'import ', 'function', 'var ', 'const ', 'let ']):
                code_lines.append(line)
        
        if code_lines:
            return "\n".join(code_lines)
        
        return content  # Return full content if no code found
    
    def _summarize_content(self, content: str) -> str:
        """Create a simple summary of content."""
        lines = content.split('\n')
        
        # Get first paragraph that looks like summary (not too short, not too long)
        for line in lines:
            line = line.strip()
            if len(line) > 50 and len(line) < 500:
                return line
        
        # Fallback: first 300 chars
        return content[:300] + "..." if len(content) > 300 else content
    
    def _extract_headings(self, content: str) -> str:
        """Extract headings from content."""
        import re
        
        # Match markdown headings
        heading_pattern = r'^(#{1,6})\s+(.+)$'
        matches = re.findall(heading_pattern, content, re.MULTILINE)
        
        if matches:
            return "\n".join([f"{'  ' * (len(h) - 1)}{text}" for h, text in matches])
        
        # Try other patterns (all caps lines, underlined, etc.)
        lines = content.split('\n')
        headings = []
        for line in lines:
            line = line.strip()
            # Heuristic: short lines, all caps, or followed by ===/---
            if line and len(line) < 100 and (line.isupper() or line.endswith(':')):
                headings.append(line)
        
        if headings:
            return "\n".join(headings)
        
        return "No headings found in document."
    
    def _read_text(self, file_path: str) -> str:
        """Read plain text files."""
        try:
            import os
            logger.info(f"Reading text file: {file_path}, exists={os.path.exists(file_path)}, size={os.path.getsize(file_path) if os.path.exists(file_path) else 0}")
            
            # Try UTF-8 first
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                logger.info(f"Successfully read {len(content)} characters")
                return content
        except UnicodeDecodeError:
            logger.warning(f"UTF-8 failed, trying latin-1 for {file_path}")
            # Fall back to latin-1
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
                logger.info(f"Successfully read {len(content)} characters with latin-1")
                return content
        except Exception as e:
            logger.error(f"Error reading text file: {e}")
            raise
    
    def _read_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2."""
        try:
            import PyPDF2
            
            text_parts = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
                
                full_text = "\n\n".join(text_parts)
                
                # Add metadata
                metadata = pdf_reader.metadata
                if metadata:
                    meta_text = "\n\n--- Document Metadata ---\n"
                    if metadata.title:
                        meta_text += f"Title: {metadata.title}\n"
                    if metadata.author:
                        meta_text += f"Author: {metadata.author}\n"
                    if metadata.subject:
                        meta_text += f"Subject: {metadata.subject}\n"
                    full_text += meta_text
                
                return full_text
                
        except ImportError:
            return "[PDF support not available. Install PyPDF2: pip install PyPDF2]"
        except Exception as e:
            return f"[Error reading PDF: {str(e)}]"
    
    def _read_word(self, file_path: str) -> str:
        """Extract text from Word documents."""
        try:
            import os
            logger.info(f"Reading Word document: {file_path}, exists={os.path.exists(file_path)}, size={os.path.getsize(file_path) if os.path.exists(file_path) else 0}")
            
            try:
                from docx import Document
            except ImportError as ie:
                logger.error(f"python-docx not installed: {ie}")
                return "[ERROR: python-docx library not installed. Run: pip install python-docx]"
            
            doc = Document(file_path)
            text_parts = []
            
            # Extract paragraphs
            para_count = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
                    para_count += 1
            
            logger.info(f"Extracted {para_count} paragraphs from Word document")
            
            # Extract tables
            table_count = 0
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n[Table]\n" + "\n".join(table_text) + "\n[/Table]")
                    table_count += 1
            
            logger.info(f"Extracted {table_count} tables from Word document")
            
            full_text = "\n\n".join(text_parts)
            logger.info(f"Total Word document content: {len(full_text)} characters")
            
            if not full_text.strip():
                return "[Document appears to be empty or contains no readable text]"
            
            return full_text
            
        except ImportError as e:
            logger.error(f"Import error for Word reading: {e}")
            return f"[ERROR: {str(e)} - Install python-docx: pip install python-docx]"
        except Exception as e:
            logger.exception(f"Error reading Word document: {e}")
            return f"[ERROR reading Word document: {str(e)}]"
    
    def _read_excel(self, file_path: str) -> str:
        """Extract text from Excel files."""
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text_parts.append(f"--- Sheet: {sheet_name} ---\n")
                text_parts.append(df.to_string(index=False))
                text_parts.append("\n")
            
            return "\n".join(text_parts)
            
        except ImportError:
            return "[Excel support not available. Install pandas and openpyxl: pip install pandas openpyxl]"
        except Exception as e:
            return f"[Error reading Excel file: {str(e)}]"
    
    def _read_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # Only add if there's content
                    text_parts.append("\n".join(slide_text))
            
            return "\n\n".join(text_parts)
            
        except ImportError:
            return "[PowerPoint support not available. Install python-pptx: pip install python-pptx]"
        except Exception as e:
            return f"[Error reading PowerPoint: {str(e)}]"
    
    def _read_image_ocr(self, file_path: str) -> str:
        """Extract text from images using OCR."""
        try:
            from PIL import Image
            import pytesseract
            
            # Open image
            image = Image.open(file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image)
            
            # Add image metadata
            width, height = image.size
            metadata = f"[Image: {width}x{height} pixels, Mode: {image.mode}]\n\n"
            
            return metadata + text if text.strip() else metadata + "[No text detected in image]"
            
        except ImportError:
            return "[OCR support not available. Install pytesseract and Pillow: pip install pytesseract Pillow]"
        except Exception as e:
            return f"[Error reading image: {str(e)}]"
    
    def _save_to_memory(self, result: Dict[str, Any]) -> bool:
        """Save document content to JARVIS memory."""
        try:
            from jarvis.memory import memory
            
            # Create a summary for conversation memory
            file_name = result.get('file_name', 'Unknown')
            file_type = result.get('file_type', '.txt')
            word_count = result.get('word_count', 0)
            summary = result.get('summary', '')[:200]
            
            conversation_summary = (
                f"Read document: {file_name} ({file_type}, {word_count} words). "
                f"Content preview: {summary}..."
            )
            
            # Save to conversations
            memory.save_conversation(
                summary=conversation_summary,
                topics=["document_reading", file_type.replace('.', '')],
                importance=2  # Higher importance for document reading
            )
            
            # Save full content as a preference for easy retrieval
            content_key = f"doc_content_{file_name}_{int(os.path.getmtime(result['file_path']))}"
            memory.save_preference(
                key=content_key,
                value=result.get('content', '')[:10000],  # Limit to 10KB
                category="documents"
            )
            
            # Save metadata as separate preference
            meta_key = f"doc_meta_{file_name}_{int(os.path.getmtime(result['file_path']))}"
            metadata = {
                "file_name": file_name,
                "file_path": result['file_path'],
                "file_type": file_type,
                "word_count": word_count,
                "summary": result.get('summary', '')
            }
            import json
            memory.save_preference(
                key=meta_key,
                value=json.dumps(metadata),
                category="documents"
            )
            
            logger.info(f"📄 Document saved to memory: {file_name}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to save document to memory: {e}")
            return False
    
    def search_in_documents(self, query: str) -> List[Dict[str, Any]]:
        """Search for content in saved documents."""
        try:
            from jarvis.memory import memory
            import json
            
            results = []
            prefs = memory.get_all_preferences(category="documents")
            
            for key, value in prefs.items():
                if key.startswith("doc_content_"):
                    if query.lower() in value.lower():
                        # Get metadata
                        meta_key = key.replace("doc_content_", "doc_meta_")
                        metadata = {}
                        if meta_key in prefs:
                            try:
                                metadata = json.loads(prefs[meta_key])
                            except:
                                pass
                        
                        # Find snippet around query
                        idx = value.lower().find(query.lower())
                        snippet_start = max(0, idx - 100)
                        snippet_end = min(len(value), idx + len(query) + 100)
                        snippet = value[snippet_start:snippet_end]
                        
                        results.append({
                            "file_name": metadata.get('file_name', 'Unknown'),
                            "file_path": metadata.get('file_path', ''),
                            "snippet": snippet,
                            "match_count": value.lower().count(query.lower())
                        })
            
            return results
            
        except Exception as e:
            logger.error(f"Error searching documents: {e}")
            return []
    
    def list_saved_documents(self) -> List[Dict[str, Any]]:
        """List all documents saved in memory."""
        try:
            from jarvis.memory import memory
            import json
            
            docs = []
            prefs = memory.get_all_preferences(category="documents")
            
            for key, value in prefs.items():
                if key.startswith("doc_meta_"):
                    try:
                        metadata = json.loads(value)
                        docs.append(metadata)
                    except:
                        pass
            
            return docs
            
        except Exception as e:
            logger.error(f"Error listing documents: {e}")
            return []
    
    def get_document_content(self, file_name: str) -> Optional[str]:
        """Retrieve full content of a saved document."""
        try:
            from jarvis.memory import memory
            
            prefs = memory.get_all_preferences(category="documents")
            
            # Find matching document
            for key, value in prefs.items():
                if key.startswith("doc_content_") and file_name in key:
                    return value
            
            return None
            
        except Exception as e:
            logger.error(f"Error retrieving document: {e}")
            return None


# Global instance
document_reader = DocumentReader()


def read_document(file_path: str, save_to_memory: bool = True) -> Dict[str, Any]:
    """Public function to read a document."""
    return document_reader.read_document(file_path, save_to_memory)


def search_documents(query: str) -> List[Dict[str, Any]]:
    """Search in saved documents."""
    return document_reader.search_in_documents(query)


def list_documents() -> List[Dict[str, Any]]:
    """List all saved documents."""
    return document_reader.list_saved_documents()


def get_document(file_name: str) -> Optional[str]:
    """Get document content by name."""
    return document_reader.get_document_content(file_name)
